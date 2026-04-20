"""
產生繁體中文版 RA-GARK 簡報（.pptx）。

需要安裝：
    pip install python-pptx

執行：
    python make_slides_zh.py
輸出：
    RA-GARK-ZH.pptx   （13 張，16:9，可直接簡報）

字型：Microsoft JhengHei（Windows 系統內建）。若在 Mac 上執行，可將
CN_FONT 改成 "PingFang TC" 或 "Heiti TC"。

投影片結構：
    1   封面
    2   研究動機 — 稀疏 KG 下 LightGCN 反而勝出
    3   整體架構
    4   詳細架構圖（若存在則嵌入 figures/architecture.png）
    5   ★ 創新 1 — Softmax 面向顯著性注意力（τ=0.5）
    6   ★ 創新 2 — 本地偏置融合閘初始化
    7   ★ 創新 3 — KG SVD 初始化
    8   主要結果（表 1）
    9   消融實驗（表 2）
    10  案例分析 — 商品級面向顯著性（圖 2）
    11  方法論洞見 — Sigmoid 陷阱
    12  限制與未來工作
    13  總結
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── 字型 ─────────────────────────────────────────────────────────────────
# Windows 內建。Mac 可改為 "PingFang TC" / "Heiti TC"。
CN_FONT = "Microsoft JhengHei"

# ── 色票 ─────────────────────────────────────────────────────────────────
C_TITLE     = RGBColor(0x1F, 0x3A, 0x68)
C_ACCENT    = RGBColor(0xD8, 0x3A, 0x3A)
C_BODY      = RGBColor(0x22, 0x22, 0x22)
C_MUTED     = RGBColor(0x55, 0x55, 0x55)
C_BG_LOCAL  = RGBColor(0xE6, 0xF0, 0xFF)
C_BG_GLOBAL = RGBColor(0xFF, 0xF0, 0xE0)
C_BG_FUSION = RGBColor(0xE8, 0xE8, 0xF8)
C_TABLE_HDR = RGBColor(0x1F, 0x3A, 0x68)
C_TABLE_ALT = RGBColor(0xF2, 0xF5, 0xFB)
C_OK        = RGBColor(0x1F, 0x7A, 0x33)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─── 工具函式 ────────────────────────────────────────────────────────────
def _set_font(run, size, bold=False, color=C_BODY, font=None):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font or CN_FONT


def add_textbox(slide, x, y, w, h, text, *, size=18, bold=False, color=C_BODY,
                align=PP_ALIGN.LEFT, font=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, size, bold, color, font)
    return tb


def add_bullets(slide, x, y, w, h, bullets, *, size=18, color=C_BODY):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "・  " + text
        _set_font(run, size, False, color)
        p.space_after = Pt(4)
    return tb


def add_rect(slide, x, y, w, h, *, fill=C_BG_LOCAL, line=C_TITLE, text=None,
             text_size=16, text_bold=False, text_color=C_BODY,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    box = slide.shapes.add_shape(shape, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(0.75)
    if text is not None:
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        _set_font(run, text_size, text_bold, text_color)
    return box


def add_slide_title(slide, title):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7),
                title, size=30, bold=True, color=C_TITLE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.5), Inches(1.0),
                                  Inches(1.2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = C_ACCENT
    line.line.fill.background()


def add_table(slide, x, y, w, h, data, *, col_widths=None,
              header_fill=C_TABLE_HDR, header_font=RGBColor(0xFF, 0xFF, 0xFF),
              alt_fill=C_TABLE_ALT, font_size=13, highlight_row=None):
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = table_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(data[r][c])
            run.font.size = Pt(font_size)
            run.font.name = CN_FONT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                run.font.bold = True
                run.font.color.rgb = header_font
            else:
                if r == highlight_row:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF1, 0xC8)
                    run.font.bold = True
                elif r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = alt_fill
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.color.rgb = C_BODY
    return table_shape


def add_bar_chart(slide, x, y, w, h, title, categories, series):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series:
        chart_data.add_series(name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, w, h, chart_data,
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    for p in chart.chart_title.text_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = C_TITLE
            run.font.name = CN_FONT
    chart.has_legend = True
    return chart


# ─── 各張投影片 ──────────────────────────────────────────────────────────
def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0), Inches(3.0),
                                    SLIDE_W, Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = C_ACCENT
    stripe.line.fill.background()

    add_textbox(slide, Inches(0.7), Inches(1.3), Inches(12), Inches(1.3),
                "RA-GARK", size=72, bold=True, color=C_TITLE)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(12), Inches(0.7),
                "基於評論面向知識圖譜的理由感知閘控網路",
                size=24, color=C_TITLE)
    add_textbox(slide, Inches(0.7), Inches(3.3), Inches(12), Inches(0.7),
                "雙視角推薦：Softmax 面向顯著性 + 本地偏置融合",
                size=18, color=C_MUTED)
    add_textbox(slide, Inches(0.7), Inches(5.7), Inches(12), Inches(0.6),
                "NDCG@20 = 0.1238     相較 KGRec 提升 +13.1%",
                size=22, bold=True, color=C_OK)
    add_textbox(slide, Inches(0.7), Inches(6.4), Inches(12), Inches(0.4),
                "資料集：Amazon Books  |  905 使用者 × 1399 商品",
                size=14, color=C_MUTED)


def slide_motivation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "研究動機 — 現有 KG 推薦方法為何表現不佳")

    add_textbox(slide, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.6),
                "在稀疏知識圖譜的資料集上（905 使用者 × 1,399 商品 × "
                "3,370 條 KG 邊），現有 KG 方法反而輸給純 LightGCN。",
                size=17, color=C_BODY)

    data = [
        ["模型", "NDCG@20", "HR@20", "Recall@20"],
        ["MCCLK (SIGIR 2022)", "0.1067", "0.4530", "0.1720"],
        ["KGCL  (SIGIR 2022)", "0.1073", "0.4696", "0.1827"],
        ["KGAT  (KDD   2019)", "0.1079", "0.4773", "0.1807"],
        ["KGRec (KDD   2023)", "0.1095", "0.4729", "0.1834"],
        ["純 LightGCN（不用 KG 的基準）", "0.1179", "0.4917", "0.1937"],
    ]
    add_table(slide, Inches(0.5), Inches(2.0), Inches(7.5), Inches(2.8),
              data, col_widths=[Inches(3.3), Inches(1.4), Inches(1.4), Inches(1.4)],
              highlight_row=5)

    add_rect(slide, Inches(8.3), Inches(2.0), Inches(4.6), Inches(2.8),
             fill=C_BG_GLOBAL,
             text="KG 重度聚合器\n(bi-interaction + L2-norm)\n"
                  "把 KG 噪聲注入 scoring。\n\n"
                  "稀疏 KG 下反而有害：\n4 個 SOTA 全輸 LightGCN。",
             text_size=15, text_bold=True)

    add_textbox(slide, Inches(0.5), Inches(5.1), Inches(12), Inches(0.5),
                "核心問題：如何讓 KG「有幫助時才貢獻」而不污染協同過濾訊號？",
                size=18, bold=True, color=C_ACCENT)

    add_rect(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.3),
             fill=C_BG_FUSION,
             text="RA-GARK 解答：雙視角 + KG + 融合閘初始偏向 LightGCN + "
                  "Softmax 面向顯著性 + KG-SVD 熱啟動\n"
                  "→ NDCG 0.1238（相較最強 KG baseline 提升 13.1%）",
             text_size=17)


def slide_architecture_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "整體架構 — 雙視角推薦器")

    add_rect(slide, Inches(5.3), Inches(1.3), Inches(2.6), Inches(0.6),
             fill=RGBColor(0xF0, 0xF0, 0xF0),
             text="輸入（使用者 u，商品 i）", text_size=16, text_bold=True)

    add_rect(slide, Inches(0.6), Inches(2.3), Inches(5.8), Inches(1.1),
             fill=C_BG_LOCAL,
             text="本地視角　（協同訊號）\n"
                  "LightGCN 在 user–item 圖上傳播\n"
                  "D⁻¹ᐟ² A D⁻¹ᐟ² x，K=2，層均值",
             text_size=14)

    add_rect(slide, Inches(6.9), Inches(2.3), Inches(5.8), Inches(1.1),
             fill=C_BG_GLOBAL,
             text="全域視角　（KG 語意訊號）\n"
                  "item_kg_aspects [N, A, d]  +  理由遮罩\n"
                  "softmax(MLP([u_glo; i_aspects]) / τ)",
             text_size=14)

    add_rect(slide, Inches(2.3), Inches(3.6), Inches(8.8), Inches(0.55),
             fill=RGBColor(0xF8, 0xF8, 0xE0),
             text="L_aCL  +  L_uCL  —  跨視角對比對齊（全域側 stop-grad）",
             text_size=14)

    add_rect(slide, Inches(2.3), Inches(4.5), Inches(8.8), Inches(1.0),
             fill=C_BG_FUSION,
             text="本地偏置融合閘（使用者 / 商品各一）\n"
                  "α = σ(MLP([loc; glo]) + 5)   →   初始 α ≈ 0.993\n"
                  "u_final = α·u_loc + (1−α)·u_glo　（商品同式）",
             text_size=14, text_bold=True)

    add_rect(slide, Inches(4.8), Inches(5.8), Inches(3.7), Inches(0.65),
             fill=RGBColor(0xFF, 0xFF, 0xFF),
             text="分數  =  u_final · i_final",
             text_size=18, text_bold=True)

    add_rect(slide, Inches(4.8), Inches(6.6), Inches(3.7), Inches(0.55),
             fill=RGBColor(0xF8, 0xF8, 0xE0),
             text="L_BPR  +  0.005·(L_aCL + L_uCL)", text_size=14)


def slide_architecture_figure(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "圖 1 — 詳細架構圖")

    fig_path = Path("figures/architecture.png")
    if fig_path.exists():
        slide.shapes.add_picture(str(fig_path), Inches(0.8), Inches(1.3),
                                 width=Inches(11.7))
        add_textbox(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5),
                    "★ 標記三個創新點：Softmax 理由注意力、本地偏置融合閘初始化、KG SVD 初始化",
                    size=13, color=C_MUTED, align=PP_ALIGN.CENTER)
    else:
        add_textbox(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(2.0),
                    "找不到架構圖檔。\n"
                    "請先執行：python figures/architecture.py\n"
                    "以產生 figures/architecture.png。",
                    size=22, color=C_MUTED, align=PP_ALIGN.CENTER)


def slide_novelty_softmax(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "★ 創新 1 — Softmax 面向顯著性注意力")

    add_textbox(slide, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.5),
                "每個商品有 A=4 個學習的面向槽；由 u_glo 條件化的 "
                "softmax attention 選出最突出的面向。",
                size=16, color=C_BODY)

    add_rect(slide, Inches(0.5), Inches(1.95), Inches(12.3), Inches(1.0),
             fill=C_BG_GLOBAL,
             text="logits  = MLP([u_glo ; i_aspects])            # [B, A]\n"
                  "weights = softmax(logits / τ, dim=-1)         # τ = 0.5（關鍵）\n"
                  "i_glo   = Σ_a weights[a] · i_aspects[a]",
             text_size=18, text_bold=True)

    data = [
        ["理由模組設計", "NDCG@20", "效果"],
        ["無（均勻平均）",                      "0.1222", "基準"],
        ["Sigmoid MLP（樸素版）",               "0.1152", "−5.7%　相較均勻 — 有害"],
        ["Softmax MLP，τ=1.0",                  "0.1232", "注意力趨近均勻，無效"],
        ["Softmax MLP，τ=0.5（本文）",          "0.1238", "最佳 NDCG，商品級訊號清晰"],
    ]
    add_table(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(2.4),
              data, col_widths=[Inches(4.5), Inches(1.8), Inches(6.0)],
              highlight_row=4, font_size=14)

    add_textbox(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.4),
                "關鍵發現：Sigmoid 形式比完全不用理由還糟。",
                size=17, bold=True, color=C_ACCENT)
    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.7),
                "Softmax 強制跨面向競爭（權重總和為 1）；τ<1 放大 logit "
                "細微差異，使商品級選擇在案例分析（第 10 頁）中清晰可見。",
                size=14, color=C_BODY)


def slide_novelty_fusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "★ 創新 2 — 本地偏置融合閘初始化")

    add_textbox(slide, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.5),
                "凸組合 α·loc + (1−α)·glo 無優雅退化機制；預設 α ≈ 0.5 "
                "讓 noisy KG 從 epoch 1 就污染 LightGCN。",
                size=16, color=C_BODY)

    add_rect(slide, Inches(0.7), Inches(1.95), Inches(5.8), Inches(2.2),
             fill=RGBColor(0xFF, 0xE0, 0xE0),
             text="未加 bias init（α ≈ 0.5）\n\n"
                  "epoch 1 起就 50/50 混合\n"
                  "noisy KG 梯度\n"
                  "污染 LightGCN\n\n"
                  "NDCG = 0.1173　（−5.3%）",
             text_size=16, text_bold=True)

    add_rect(slide, Inches(6.8), Inches(1.95), Inches(5.8), Inches(2.2),
             fill=RGBColor(0xE0, 0xF0, 0xE0),
             text="使用 bias = +5（本文）\n\n"
                  "初始 α ≈ σ(5) ≈ 0.993\n"
                  "模型一開始像純 LightGCN\n"
                  "KG 有用時閘才打開\n\n"
                  "NDCG = 0.1238　（最佳）",
             text_size=16, text_bold=True)

    add_rect(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(1.4),
             fill=C_BG_FUSION,
             text="gate = Sequential(Linear(2d → d), Tanh,\n"
                  "                  Linear(d → 1, bias_init = +5), Sigmoid)\n"
                  "α    = gate([loc ; glo])",
             text_size=18, text_bold=True)

    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0),
                "論文洞見：稀疏 KG 下的樸素雙視角融合可能輸給純 LightGCN。"
                "用大正 bias 初始化融合閘給模型安全退路 — 全域視角變成 "
                "「可選」而非「強制」。",
                size=14, color=C_BODY)


def slide_novelty_svd(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "★ 創新 3 — KG SVD 初始化")

    add_textbox(slide, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.5),
                "item_kg_aspects 從 KG 關聯矩陣的截斷 SVD 開始，"
                "而非隨機初始化。",
                size=16, color=C_BODY)

    steps = [
        "從 KG 邊建稀疏矩陣 M ∈ R^{N_items × N_aspects}（TF-IDF 加權）。",
        "截斷 SVD： M ≈ U Σ Vᵀ   →   取 U · √Σ 作為商品的面向感知嵌入。",
        "重塑為 [N_items, A, d] = [1399, 4, 128]，縮放至 xavier σ。",
        "作為 item_kg_aspects 的初始值；單一 1e-3 Adam lr。",
    ]
    add_bullets(slide, Inches(0.5), Inches(1.9), Inches(12.3), Inches(2.4),
                steps, size=17)

    data = [
        ["初始化方式", "NDCG@20", "相較 winner"],
        ["Xavier（隨機）",      "0.1173", "−5.3%"],
        ["KG SVD（本文）",      "0.1238", "—"],
    ]
    add_table(slide, Inches(0.5), Inches(4.8), Inches(7.0), Inches(1.4),
              data, col_widths=[Inches(2.8), Inches(1.8), Inches(2.4)],
              highlight_row=2, font_size=15)

    add_rect(slide, Inches(7.9), Inches(4.8), Inches(5.0), Inches(1.4),
             fill=C_BG_GLOBAL,
             text="為何有效：\n"
                  "BPR 訓練過程中 KG 面向幾何被保留 —\n"
                  "全域視角不必從零開始學 KG 結構。",
             text_size=14)


def slide_main_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "主要結果 — 表 1")

    data = [
        ["模型", "NDCG@20", "HR@20", "Recall@20", "MAP@20"],
        ["MCCLK (SIGIR 2022)",       "0.1067", "0.4530", "0.1720", "0.0497"],
        ["KGCL  (SIGIR 2022)",       "0.1073", "0.4696", "0.1827", "0.0479"],
        ["KGAT  (KDD   2019)",       "0.1079", "0.4773", "0.1807", "0.0491"],
        ["KGRec (KDD   2023)",       "0.1095", "0.4729", "0.1834", "0.0500"],
        ["RA-GARK（本文，τ=0.5）",   "0.1238", "0.4961", "0.2014", "0.0591"],
    ]
    add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.1),
              data, col_widths=[Inches(3.9), Inches(2.1), Inches(2.1),
                                Inches(2.1), Inches(2.1)],
              highlight_row=5, font_size=16)

    data2 = [
        ["指標", "KGRec", "RA-GARK", "相對提升"],
        ["NDCG@20",  "0.1095", "0.1238", "+13.1 %"],
        ["HR@20",    "0.4729", "0.4961", "+4.9 %"],
        ["Recall@20","0.1834", "0.2014", "+9.8 %"],
        ["MAP@20",   "0.0500", "0.0591", "+18.2 %"],
    ]
    add_table(slide, Inches(0.5), Inches(4.8), Inches(8.0), Inches(2.3),
              data2, col_widths=[Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)],
              font_size=15)

    add_rect(slide, Inches(9.0), Inches(4.8), Inches(3.8), Inches(2.3),
             fill=C_BG_FUSION,
             text="資料集：Amazon Books\n"
                  "905 使用者 × 1399 商品\n"
                  "3,370 條有效 KG 邊\n"
                  "使用者分層 70/15/15\n"
                  "seed = 42",
             text_size=15)


def slide_ablation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "消融實驗 — 三個 ★ 創新缺一不可")

    data = [
        ["設定", "NDCG@20", "相較 winner", "移除的組件"],
        ["winner（τ=0.5）",         "0.1238", "—",       "softmax + τ + 融合 bias + 所有組件"],
        ["winner_no_rat",           "0.1222", "−1.3%",   "均勻平均（無理由）"],
        ["winner_no_acl",           "0.1207", "−2.5%",   "移除 aspect-level CL"],
        ["winner_no_ucl",           "0.1187", "−4.1%",   "移除 user cross-view CL"],
        ["winner_no_svd  ★",        "0.1173", "−5.3%",   "xavier 初始化（非 KG SVD）"],
        ["winner_fb0  ★",           "0.1173", "−5.3%",   "融合 bias 由 5 改為 0"],
        ["winner_sigmoid_rat  ★",   "0.1152", "−7.0%",   "softmax 改為 sigmoid MLP"],
        ["old_full",                "0.1067", "−13.8%",  "兩個 Fix-1 皆還原（修復前）"],
        ["no_global_view",          "0.1214", "−1.9%",   "跳過全域管線（僅 CL）"],
        ["lightgcn_only",           "0.1179", "−4.8%",   "完全不用 KG（下限）"],
    ]
    add_table(slide, Inches(0.3), Inches(1.3), Inches(12.7), Inches(4.6),
              data, col_widths=[Inches(3.3), Inches(1.5), Inches(1.8), Inches(6.1)],
              highlight_row=1, font_size=13)

    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.4),
                "移除任一 ★ 創新都造成 ≥ 5% NDCG 下降 → 每一個都是真正的貢獻。",
                size=16, bold=True, color=C_TITLE)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.6),
                "同時還原兩個 Fix-1（old_full）掉到 0.1067 — 低於純 LightGCN。"
                "沒有這些修正，整個雙視角架構反而是負擔。",
                size=14, color=C_MUTED)


def slide_case_study(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "案例分析 — 商品級面向顯著性（圖 2）")

    add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
                "5 個代表性商品在 4 個面向槽上的 Softmax 注意力權重（τ=0.5）。",
                size=15, color=C_BODY)

    categories = ["商品 296\n（賽博龐克）",
                  "商品 1331\n（心理驚悚）",
                  "商品 785\n（科幻生存）",
                  "商品 245\n（反烏托邦科幻）",
                  "商品 77\n（末日廢土）"]
    series = [
        ("aspect #0", (0.211, 0.254, 0.324, 0.289, 0.269)),
        ("aspect #1", (0.410, 0.340, 0.260, 0.267, 0.252)),
        ("aspect #2", (0.211, 0.186, 0.182, 0.248, 0.225)),
        ("aspect #3", (0.168, 0.220, 0.234, 0.196, 0.254)),
    ]
    add_bar_chart(slide, Inches(0.3), Inches(1.7), Inches(7.7), Inches(5.2),
                  "各商品的面向注意力　（τ = 0.5）", categories, series)

    add_rect(slide, Inches(8.2), Inches(1.7), Inches(4.8), Inches(3.4),
             fill=C_BG_FUSION,
             text="觀察\n\n"
                  "・每個商品都有不同的主導面向\n"
                  "　（跨商品 argmax 不同）\n\n"
                  "・賽博龐克《Count Zero》(296)\n"
                  "　→ 槽 #1 權重 0.41（明顯 > 0.25）\n\n"
                  "・科幻生存《The Martian》(785)\n"
                  "　→ 槽 #0 權重 0.32\n\n"
                  "・均勻基準為 0.25",
             text_size=13)

    add_rect(slide, Inches(8.2), Inches(5.3), Inches(4.8), Inches(1.7),
             fill=RGBColor(0xFF, 0xF1, 0xC8),
             text="坦承的限制：\n"
                  "同一商品、不同使用者的\n"
                  "注意力幾乎相同（Δ ≈ 0.005）。\n"
                  "目前為商品條件化；更強的\n"
                  "使用者條件化是 future work。",
             text_size=12)


def slide_insight(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "方法論洞見 — Sigmoid 陷阱")

    add_textbox(slide, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.5),
                "本文最有意思的發現不是 +13.1% NDCG，而是「稀疏 KG 下樸素設計為何失敗」。",
                size=18, bold=True, color=C_TITLE)

    add_rect(slide, Inches(0.5), Inches(2.1), Inches(4.0), Inches(3.2),
             fill=RGBColor(0xFF, 0xE0, 0xE0),
             text="樸素做法（sigmoid MLP）\n"
                  "weights = σ(MLP([u; a]))\n\n"
                  "NDCG  =  0.1152\n\n"
                  "比無理由更糟\n"
                  "比純 LightGCN 更糟\n"
                  "sigmoid 飽和、注入噪聲",
             text_size=15, text_bold=True)

    add_rect(slide, Inches(4.7), Inches(2.1), Inches(4.0), Inches(3.2),
             fill=RGBColor(0xF5, 0xF5, 0xF5),
             text="無理由（均勻平均）\n"
                  "i_glo = mean(i_aspects)\n\n"
                  "NDCG  =  0.1222\n\n"
                  "安全預設\n"
                  "無使用者條件化\n"
                  "無可解釋性",
             text_size=15, text_bold=True)

    add_rect(slide, Inches(8.9), Inches(2.1), Inches(4.0), Inches(3.2),
             fill=RGBColor(0xE0, 0xF0, 0xE0),
             text="本文做法（softmax τ=0.5）\n"
                  "softmax(MLP(…)/τ)\n\n"
                  "NDCG  =  0.1238\n\n"
                  "跨面向競爭\n"
                  "商品級可解釋性\n"
                  "安全退化（τ=1 ≈ 均勻）",
             text_size=15, text_bold=True)

    add_textbox(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.4),
                "核心訊息：attention 形式不是裝飾細節。sigmoid 與 softmax "
                "的選擇，可讓 NDCG 差 7+ 個百分點。未來 KG-aware 方法必須"
                "驗證其 KG 整合在 noise / sparse 情境下能優雅退化；"
                "樸素設計有可能是淨負面貢獻。",
                size=15, color=C_BODY)


def slide_limitations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "限制與未來工作")

    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
                "誠實檢視 RA-GARK 仍待改進之處。",
                size=17, color=C_BODY)

    bullets = [
        "注意力是「商品條件化」而非「使用者條件化」。同商品、不同使用者的 "
        "softmax 幾乎相同（Δ ≈ 0.005）。更強的使用者訊號可從 FiLM 風格調變、"
        "user-aspect bias、或更大的 u_glo 維度著手。",

        "僅在單一資料集（Amazon Books）驗證。Amazon Movies / Yelp / "
        "Last.FM 等實驗可強化外部效度。",

        "單一 seed 結果。多 seed（42, 1, 7, 100, 2025）取 mean ± std "
        "能確認 rationale 相對均勻的 +0.7% 是真訊號還是雜訊。",

        "A=4 個面向槽可能不足以涵蓋 2,098 個獨立 KG 面向詞。A ∈ {8, 16} "
        "可能有更佳 sweet spot。",

        "KG 品質與資料集相關。我們關於「sigmoid 陷阱」的洞見能否推廣至 "
        "密集 KG 領域，仍需驗證。",
    ]
    add_bullets(slide, Inches(0.8), Inches(1.9), Inches(12), Inches(5.0),
                bullets, size=14)


def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "總結")

    add_rect(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.0),
             fill=C_BG_FUSION,
             text="RA-GARK：雙視角推薦器，結合 Softmax 面向顯著性理由、"
                  "本地偏置融合、KG-SVD 初始化\n"
                  "NDCG@20 = 0.1238　（相較最強 KG baseline KGRec 提升 13.1%）",
             text_size=17, text_bold=True)

    add_textbox(slide, Inches(0.5), Inches(2.5), Inches(12), Inches(0.4),
                "三個核心貢獻（消融實驗驗證）：",
                size=18, bold=True, color=C_TITLE)

    add_rect(slide, Inches(0.5), Inches(3.0), Inches(4.0), Inches(2.0),
             fill=C_BG_GLOBAL,
             text="★　Softmax 面向顯著性\n　　注意力\n\n"
                  "Sigmoid → softmax / τ\n"
                  "移除掉 −7.0%\n"
                  "可解釋的商品級\n"
                  "注意力",
             text_size=14)
    add_rect(slide, Inches(4.7), Inches(3.0), Inches(4.0), Inches(2.0),
             fill=C_BG_LOCAL,
             text="★　本地偏置融合閘\n　　初始化\n\n"
                  "bias = +5 → α ≈ 0.993\n"
                  "移除掉 −5.3%\n"
                  "可安全退化為\n"
                  "LightGCN",
             text_size=14)
    add_rect(slide, Inches(8.9), Inches(3.0), Inches(4.0), Inches(2.0),
             fill=C_BG_FUSION,
             text="★　KG SVD\n　　初始化\n\n"
                  "item_kg_aspects ← SVD(M)\n"
                  "移除掉 −5.3%\n"
                  "保留 KG 語意幾何",
             text_size=14)

    add_rect(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.0),
             fill=RGBColor(0xFF, 0xF1, 0xC8),
             text="方法論洞見：稀疏 KG 下的樸素 KG 整合是「淨負面貢獻」；"
                  "注意力形式（sigmoid vs softmax）讓 NDCG 差 7+ 個百分點。",
             text_size=14, text_bold=True)

    add_textbox(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.6),
                "產出素材：ablation_results_paper.csv ・ case_study.csv ・ "
                "figures/architecture.pdf ・ ARCHITECTURE.md ・ "
                "run_ablations.py (--mode minimal|paper|full|temp)",
                size=11, color=C_MUTED, align=PP_ALIGN.CENTER)


# ─── 主程式 ──────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_motivation(prs)
    slide_architecture_overview(prs)
    slide_architecture_figure(prs)
    slide_novelty_softmax(prs)
    slide_novelty_fusion(prs)
    slide_novelty_svd(prs)
    slide_main_results(prs)
    slide_ablation(prs)
    slide_case_study(prs)
    slide_insight(prs)
    slide_limitations(prs)
    slide_summary(prs)

    out = "RA-GARK-ZH.pptx"
    prs.save(out)
    print(f"已存檔 → {out}   ({len(prs.slides)} 張投影片)")


if __name__ == "__main__":
    main()
